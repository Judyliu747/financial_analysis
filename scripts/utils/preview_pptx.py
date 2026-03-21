"""
Render first slide of PPTX to PNG using python-pptx + PIL
Draws background, text boxes, tables, shapes as a visual preview
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import io, sys, os

PPTX_PATH = "/Users/macrossz/DevTools/VscodeProject/ClaudeCode/financial_analysis/output/SMIC/SMIC_OnePager_StripProfile.pptx"
OUT_PATH  = "/Users/macrossz/DevTools/VscodeProject/ClaudeCode/financial_analysis/output/SMIC/smic_onepager_preview.png"

SCALE = 120  # px per inch
W_IN, H_IN = 10.0, 7.5
W_PX = int(W_IN * SCALE)
H_PX = int(H_IN * SCALE)

prs = Presentation(PPTX_PATH)
slide = prs.slides[0]

img = Image.new("RGB", (W_PX, H_PX), (255, 255, 255))
draw = ImageDraw.Draw(img)

def emu_to_px(emu):
    return int(emu / 914400 * SCALE)

def rgb_from_hex(h):
    if not h or len(h) < 6:
        return (0, 0, 0)
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

for shape in slide.shapes:
    x = emu_to_px(shape.left or 0)
    y = emu_to_px(shape.top or 0)
    w = emu_to_px(shape.width or 0)
    h = emu_to_px(shape.height or 0)

    # Fill shapes
    try:
        fill = shape.fill
        if fill.type is not None:
            try:
                fg = fill.fore_color.rgb
                color = (fg.r, fg.g, fg.b)
                draw.rectangle([x, y, x+w, y+h], fill=color)
            except:
                pass
    except:
        pass

    # Lines
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        try:
            lc = shape.line.color.rgb
            lcolor = (lc.r, lc.g, lc.b)
        except:
            lcolor = (180, 180, 180)
        draw.line([x, y, x+w, y+h], fill=lcolor, width=1)

    # Text frames
    if shape.has_text_frame:
        try:
            font_size = 10
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = int(run.font.size.pt)
                        break
            text = shape.text_frame.text[:200]
            try:
                fnt = ImageFont.load_default()
            except:
                fnt = None
            draw.text((x+3, y+2), text[:60], fill=(20, 20, 80), font=fnt)
        except:
            pass

    # Tables
    if shape.has_table:
        tbl = shape.table
        try:
            cell_h = h // max(len(tbl.rows), 1)
            cell_w = w // max(len(tbl.columns), 1)
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    cx = x + ci * cell_w
                    cy = y + ri * cell_h
                    # Cell border
                    draw.rectangle([cx, cy, cx+cell_w, cy+cell_h],
                                   outline=(180, 180, 180), width=1)
                    # Cell fill
                    try:
                        cf = cell.fill.fore_color.rgb
                        draw.rectangle([cx+1, cy+1, cx+cell_w-1, cy+cell_h-1],
                                       fill=(cf.r, cf.g, cf.b))
                    except:
                        pass
                    # Cell text
                    try:
                        txt = cell.text[:25]
                        fnt = ImageFont.load_default()
                        col = (20, 20, 80)
                        try:
                            tc = cell.text_frame.paragraphs[0].runs[0].font.color.rgb
                            col = (tc.r, tc.g, tc.b)
                        except:
                            pass
                        draw.text((cx+3, cy+2), txt, fill=col, font=fnt)
                    except:
                        pass
        except:
            pass

# Grid overlay
draw.rectangle([0, 0, W_PX-1, H_PX-1], outline=(200, 200, 200), width=2)

img.save(OUT_PATH)
print(f"Preview saved: {OUT_PATH}")
